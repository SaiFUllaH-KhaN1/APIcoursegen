# from langchain_openai import ChatOpenAI
import PROMPTS as PROMPTS
import PROMPTS_WITHOUT_FILE as PROMPTS_WITHOUT_FILE
import logging
from langchain.chains import LLMChain
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
import base64
from langchain.schema.messages import HumanMessage, SystemMessage
# from PyPDF2 import PdfReader
from pypdf import PdfReader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
import os
import shutil
import json
from langchain.schema.document import Document
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders import UnstructuredExcelLoader, UnstructuredWordDocumentLoader, UnstructuredPowerPointLoader, TextLoader
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
from pydantic import BaseModel, Field
from typing import List, Dict, Any, Optional
from docx2python import docx2python
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import shutil
import re
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from PIL import Image
from langchain.utils.math import cosine_similarity
import urllib
import urllib.request
from urllib.parse import urlparse, urljoin
# from cairosvg import svg2png
from bs4 import BeautifulSoup
import imagehash


import traceback
import fitz
# import subprocess # Used for the libreoffice command for ppt and doc support (NOT pptx AND docx) 

# Logging Declaration
log_format = '%(asctime)s - %(levelname)s - %(message)s'
logger = logging
logger.basicConfig(level= logging.DEBUG, format= log_format)

def RAG(file_content,embeddings,file,session_var, temp_path_audio,filename, extension, language, temp_pdf_file):
    logger.info(f"file is: {file}",)
    
    filename_without_extension = filename.rsplit('.', 1)[0].lower()
    if f"extracted_content{session_var}" not in filename_without_extension:
        output_path_byfile = f"./imagefolder_{session_var}/images_{session_var}_{filename_without_extension}"
        if not os.path.exists(output_path_byfile):
            os.makedirs(output_path_byfile, exist_ok=True) 
    logger.info(f"Extension is: {extension}",)
    raw_text = ''
    texts = '' # for pdf image path only!
    if extension=="pdf":

        # For Image processing
        if f'extracted_content{session_var}.pdf' not in filename:
            fitz_pdf_reader = fitz.open(temp_pdf_file)
            pgcount=0
            for page_num in range(len(fitz_pdf_reader)):
                page = fitz_pdf_reader.load_page(page_num)
                pgcount += 1
                imgcount = 1
                text_instant = page.get_text()
                text_instant = f"\nThe Content of PageNumber {pgcount} of file name {filename_without_extension} is:\n{text_instant}.\nEnd of PageNumber {pgcount} of file name {filename_without_extension}\n"
                try:
                    images = page.get_images(full=True)

                    # This checks if images are more than 10 on one page then they might need flattening
                    # as being non-cohesive. So, we convert the whole page to an image for summarizing
                    # in the vision model later down the line.
                    if len(images)>10: 
                        pix = page.get_pixmap(dpi=130)
                        base_name_pixmap = f"Pixmapped FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}"
                        pix.save(f"{output_path_byfile}/{base_name_pixmap}.png")

                    else:
                        for img_index, img in enumerate(images, start=1):
                            xref = img[0]
                            base_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}"
                            try:
                                base_image = fitz_pdf_reader.extract_image(xref)
                            except Exception as e:
                                logger.info(f"Error base image : {e}. Lets resolve via cmyk dealing")
                                logger.info(traceback.format_exc())
                                pix = fitz.Pixmap(fitz_pdf_reader, xref) # create a Pixmap
                                if pix.n > 4: # CMYK: convert to RGB first
                                    pix = fitz.Pixmap(fitz.csRGB, pix)
                                png_filename = f"{base_name}.png"
                                png_path = os.path.join(output_path_byfile, png_filename)
                                pix.writePNG(png_path)  # Write image content to PNG
                                pix = None
                                logger.info(f"Error base image Resolved as far as cmyk is concerned")
                                continue
                            
                            image_bytes = base_image["image"]
                            extension = base_image["ext"]
                            logger.info(f"Img extension is: {extension} and img base name is: {base_name}")
                            if extension == ".jp2":
                                logger.info(f"Checking Image Extension {extension}",)
                                image_path = os.path.join(output_path_byfile, base_name)  # Construct the full output path
                                imgcount += 1
                                with open(image_path, "wb") as fp:
                                    fp.write(image_bytes)
                                with Image.open(image_path) as im:
                                    new_image_name = base_name + ".png"  # New image name for PNG
                                    new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
                                    im.save(new_image_path)  # Save the image as PNG
                                    os.remove(image_path)
                            else:
                                image_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}.{extension}"  # Construct new file name with count
                                image_path = os.path.join(output_path_byfile, image_name)  # Construct the full output path
                                imgcount += 1
                                with open(image_path, "wb") as fp:
                                    fp.write(image_bytes)
                                with Image.open(image_path) as im:
                                    if im.mode in ["P", "PA"]:
                                        logger.info(f"Image of P or PA Mode detected:{im.mode}",)
                                        im = im.convert("RGBA")  # Convert palette-based images to RGBA
                                        new_image_name = base_name + ".png"  # New image name for PNG
                                        new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
                                        im.save(new_image_path)  # Save the image as PNG
                                        os.remove(image_path)
                                    else:
                                        pass

                    #logger.info("filler print") # suggests that the try block has executed
                except Exception as e:
                    logger.info(f"Error processing image {base_name}: {e}")
                    logger.info(traceback.format_exc())
                    pass
                
                if text_instant:
                    texts += text_instant
                    
            # at this point we got text and images extracted, stored in ./images folder, lets summarize images

            # Get image summaries
            # image_elements = []
            # image_summaries = []
            # image_path = output_path_byfile

            # def encode_image(image_path):
            #     with open(image_path, "rb") as f:
            #         return base64.b64encode(f.read()).decode('utf-8')

            # def summarize_image(encoded_image, basename):
            #     prompt = [
            #         SystemMessage(content="You are a bot that is good at analyzing images."),
            #         HumanMessage(content=[
            #             {
            #                 "type": "text",
            #                 "text": f"Describe the contents of this image. Tell what FileName, PageNumber and ImageNumber of this image is by seeing this information: {basename}. Your output should look like this: 'This image that belongs to FileName: ..., PageNumber: ..., ImageNumber: .... In this Image ...'"
            #             },
            #             {
            #                 "type": "image_url",
            #                 "image_url": {
            #                     "url": f"data:image/jpeg;base64,{encoded_image}"
            #                 },
            #             },
            #         ])
            #     ]
            #     response = ChatOpenAI(model="gpt-4o", max_tokens=128).invoke(prompt)
            #     return response.content

            # for i in os.listdir(output_path_byfile):
            #     if i.endswith(('.png', '.jpg', '.jpeg')):
            #         image_path = os.path.join(output_path_byfile, i)
            #         basename = os.path.basename(image_path) 
            #         print(os.path.basename(image_path))
            #         encoded_image = encode_image(image_path)
            #         image_elements.append(encoded_image)
            #         summary = summarize_image(encoded_image,basename)
            #         image_summaries.append(summary)

            # print("Image Summary is::",image_summaries)

            # text_merged = texts + "\n" + str(image_summaries)
            # print(text_merged) # Will be used onwards to create the faiss_index text database
            
        else:
            logger.info("We are processing for url or youtube's extracted_content.pdf")
            # Without the Image Processing
            doc_reader = PdfReader(file_content)

            for i, page in enumerate(doc_reader.pages):
                text = page.extract_text()

                if text:
                    raw_text += text

    elif extension=="csv":
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        try:
            file.save(temp_path)
            loader = CSVLoader(file_path=temp_path)
            data = loader.load()
            raw_text = raw_text.join(document.page_content + '\n\n' for document in data)
            os.remove(temp_path)
        except Exception as e:
            logger.info(f"Error in {filename}: {e}")
            logger.info(traceback.format_exc())
            os.remove(temp_path)

    elif extension=="xlsx" or extension=="xls":
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        try:
            file.save(temp_path)
            loader = UnstructuredExcelLoader(temp_path)
            data = loader.load()
            raw_text = raw_text.join(document.page_content + '\n\n' for document in data)
            os.remove(temp_path)
        except Exception as e:
            logger.info(f"Error in {filename}: {e}")
            logger.info(traceback.format_exc())
            os.remove(temp_path)


    elif extension=="docx":
        # temp_path = os.path.join(filename)
        # file.seek(0)
        # file.save(temp_path)
        # loader = UnstructuredWordDocumentLoader(temp_path)
        # data = loader.load()
        # raw_text = raw_text.join(document.page_content for document in data)
        # os.remove(temp_path)
        def extract_and_rename_images(docx_path, output_dir):
            # Extract the contents of the DOCX file
            content = docx2python(docx_path, extract_image=True,  image_folder=output_dir)
            logger.info(content)
            # Flatten the list of lists containing the images
            images = content.images
            image_info = []
            for image_name in images.keys():
                base_name, ext = os.path.splitext(image_name)
                logger.info(f"Base Name: {base_name}, Extension: {ext}")
                image_info.append((base_name, ext))
            # Rename images based on their location
            image_count = 1
            i = 0
            texts = ""
            for section_index, section in enumerate(content.body):
                for row_index, row in enumerate(section):
                    for cell_index, cell in enumerate(row):
                        for paragraph_index, paragraph in enumerate(cell):
                            if '----media/' in paragraph:
                                # Extract the image filename from the placeholder
                                start_index = paragraph.find('----media/') + len('----media/')
                                end_index = paragraph.find('----', start_index)
                                image_filename = paragraph[start_index:end_index]

                                # Construct the original image path
                                original_image_path = os.path.join(output_dir, image_filename)
                                # base_name = os.path.splitext(os.path.basename(image_file_object.name))[0]  # Get the base file name without extension
                                # extension = os.path.splitext(image_file_object.name)[1]  # Get the file extension
                                # Create a new image name based on its location
                                base_name,ext = image_info[i]
                                new_image_name = f"FileName {filename_without_extension} PageNumber Null ImageNumber {image_count}{ext}"
                                new_image_path = os.path.join(output_dir, new_image_name)

                                # Rename the image
                                if os.path.exists(original_image_path):
                                    shutil.move(original_image_path, new_image_path)

                                    # Update the placeholder in the paragraph ----media/ImageNumber (\d+) PageNumber Null of FileName (.+)----
                                    new_placeholder = f"----media/ImageNumber {image_count} PageNumber Null of FileName {filename_without_extension}----"
                                    paragraph = paragraph.replace(f"----media/{image_filename}----", new_placeholder)
                                    cell[paragraph_index] = paragraph
                                    image_count += 1
                                    i += 1
                            # print(paragraph)
                            texts += paragraph + "\n"
            return texts

        texts = extract_and_rename_images(file_content, output_path_byfile)
        logger.info(f"texts is:::{texts}",)

    elif extension=="pptx":
        # temp_path = os.path.join(filename)
        # file.seek(0)
        # file.save(temp_path)
        # loader = UnstructuredPowerPointLoader(temp_path)
        # data = loader.load()
        # raw_text = raw_text.join(document.page_content for document in data)
        # os.remove(temp_path)
        def iter_picture_shapes(prs):
            slide_number = 1
            # image_number = 1
            # img_names = []
            for slide in prs.slides:
                image_number = 1
                logger.info(f"slide {slide}",)
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            for s in shape.shapes:
                                if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    image = shape.image
                                    logger.info(f"image_number {image_number}",)
                                    image_filename = f'FileName {filename_without_extension} SlideNumber {slide_number} ImageNumber {image_number}.{image.ext}'
                                    # img = f'SlideNumber:{slide_number} of FileName:{filename_without_extension}-ImageNumber {image_number}'
                                    image_number += 1
                                    logger.info(image_filename)
                                    # img_names.append(img)

                                    image_path = os.path.join(output_path_byfile, image_filename)
                                    with open(image_path, "wb") as fp:
                                        fp.write(image.blob)
                    except Exception as e:
                        logger.info(f"Error processing image of slide {slide_number} and img no: {image_number}: {e}")
                        logger.info(traceback.format_exc())
                        pass 
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            logger.info(f"image_number {image_number}",)
                            image_filename = f'FileName {filename_without_extension} SlideNumber {slide_number} ImageNumber {image_number}.{image.ext}'
                            # img = f'SlideNumber:{slide_number} of FileName:{filename_without_extension} with ImageNumber:{image_number}\n'
                            image_number += 1
                            logger.info(image_filename)
                            # img_names.append(img)
                            image_path = os.path.join(output_path_byfile, image_filename)
                            with open(image_path, "wb") as fp:
                                fp.write(image.blob)
                    except Exception as e:
                        logger.info(f"Error processing image of slide {slide_number} and img no: {image_number}: {e}")
                        logger.info(traceback.format_exc())
                        pass 

                slide_number += 1  
               
        try:
            iter_picture_shapes(Presentation(file_content))

            # langchain unstructuredworddoc method
            temp_path = os.path.join(f"{filename}{session_var}")
            file.seek(0)
            file.save(temp_path)
            loader = UnstructuredPowerPointLoader(temp_path,mode='elements')
            data = loader.load()
            logger.info(f"data:\n{data}",)

            # Step 1: Collect content for each page number
            page_contents = {}
            for doc in data:
                # Access the metadata and page content correctly
                page_number = doc.metadata.get('page_number')
                if page_number is not None:
                    if page_number not in page_contents:
                        page_contents[page_number] = []
                    page_contents[page_number].append(doc.page_content)

            # Step 2: Combine the content for each page number
            # combined_page_contents = [{'page_number': page,'filename': filename, 'page_content': ' '.join(contents)} for page, contents in page_contents.items()]
            # pattern_this_end_pptx = r"End of SlideNumber (\d+) with Filename (.+?) ----" 
            combined_page_contents = [
                {
                    'page_content': f"\nSlideNumber {page} FileName {filename_without_extension} ---\n{' '.join(contents)} End of SlideNumber {page} with Filename {filename_without_extension} ---"
                }
                for page, contents in page_contents.items()
            ]

            texts = str(combined_page_contents)
            logger.info(texts)
            os.remove(temp_path)

        except Exception as e:
            logger.info(f"Error in {filename}: {e}")
            logger.info(traceback.format_exc())
            os.remove(temp_path)

########################Below is the ppt and doc compatibility code(use libreoffice in docker to work)########################
    # elif extension=="ppt" or extension=="doc":
    #     logger.info(f"PPT or DOC file name is ::{filename}")
    #     temp_path = os.path.join(f"{session_var}{filename}")
    #     file.seek(0)
    #     file.save(temp_path)

    #     def convert_pptordoc_to_pdf(input_path, output_dir):
    #         # Construct the command to convert pptx to pdf
    #         # ALERT--ALERT 
    #         # for ThingLink github, use only "soffice" in the command
    #         command = [
    #             'E:\Softwares Installed\LIBREOFFICE\program\soffice.exe',
    #             '--headless',
    #             '--convert-to',
    #             'pdf',
    #             '--outdir',
    #             output_dir,
    #             input_path
    #         ]
    #         # Execute the command
    #         subprocess.run(command, check=True)
    #         print(f'Converted {input_path} to PDF and saved in {output_dir}')
        
    #     # Convert the .pptx file to PDF
    #     convert_pptordoc_to_pdf(temp_path, f"pdf_dir{session_var}")
    #     os.remove(temp_path)

    #     # Processing created PDF from ppt
    #     temp_pdf_file = os.path.join(f"pdf_dir{session_var}", f"{session_var}{filename_without_extension}.pdf")
    #     fitz_pdf_reader = fitz.open(temp_pdf_file)
    #     pgcount=0
    #     for page_num in range(len(fitz_pdf_reader)):
    #         page = fitz_pdf_reader.load_page(page_num)
    #         pgcount += 1
    #         imgcount = 1
    #         text_instant = page.get_text()
    #         text_instant = f"\nThe Content of PageNumber {pgcount} of file name {filename_without_extension} is:\n{text_instant}.\nEnd of PageNumber {pgcount} of file name {filename_without_extension}\n"
    #         try:
    #             images = page.get_images(full=True)
    #             for img_index, img in enumerate(images, start=1):
    #                 xref = img[0]
    #                 base_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}"
    #                 try:
    #                     base_image = fitz_pdf_reader.extract_image(xref)
    #                 except Exception as e:
    #                     logger.info(f"Error base image : {e}. Lets resolve via cmyk dealing")
    #                     logger.info(traceback.format_exc())
    #                     pix = fitz.Pixmap(fitz_pdf_reader, xref) # create a Pixmap
    #                     if pix.n > 4: # CMYK: convert to RGB first
    #                         pix = fitz.Pixmap(fitz.csRGB, pix)
    #                     png_filename = f"{base_name}.png"
    #                     png_path = os.path.join(output_path_byfile, png_filename)
    #                     pix.writePNG(png_path)  # Write image content to PNG
    #                     pix = None
    #                     logger.info(f"Error base image Resolved as far as cmyk is concerned")
    #                     continue
                    
    #                 image_bytes = base_image["image"]
    #                 extension = base_image["ext"]
    #                 logger.info(f"Img extension is: {extension} and img base name is: {base_name}")
    #                 if extension == ".jp2":
    #                     logger.info(f"Checking Image Extension {extension}",)
    #                     image_path = os.path.join(output_path_byfile, base_name)  # Construct the full output path
    #                     imgcount += 1
    #                     with open(image_path, "wb") as fp:
    #                         fp.write(image_bytes)
    #                     with Image.open(image_path) as im:
    #                         new_image_name = base_name + ".png"  # New image name for PNG
    #                         new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
    #                         im.save(new_image_path)  # Save the image as PNG
    #                         os.remove(image_path)
    #                 else:
    #                     image_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}.{extension}"  # Construct new file name with count
    #                     image_path = os.path.join(output_path_byfile, image_name)  # Construct the full output path
    #                     imgcount += 1
    #                     with open(image_path, "wb") as fp:
    #                         fp.write(image_bytes)
    #                     with Image.open(image_path) as im:
    #                         if im.mode in ["P", "PA"]:
    #                             logger.info(f"Image of P or PA Mode detected:{im.mode}",)
    #                             im = im.convert("RGBA")  # Convert palette-based images to RGBA
    #                             new_image_name = base_name + ".png"  # New image name for PNG
    #                             new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
    #                             im.save(new_image_path)  # Save the image as PNG
    #                             os.remove(image_path)
    #                         else:
    #                             pass

    #             #logger.info("filler print") # suggests that the try block has executed
    #         except Exception as e:
    #             logger.info(f"Error processing image {base_name}: {e}")
    #             logger.info(traceback.format_exc())
    #             pass
            
    #         if text_instant:
    #             texts += text_instant
########################Above is the ppt and doc compatibility########################

    elif extension=="txt":
        logger.info(f"Text file name is ::{filename}")
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        try:
            file.save(temp_path)
            loader = TextLoader(temp_path)
            data = loader.load()
            raw_text = raw_text.join(document.page_content for document in data)
            os.remove(temp_path)
        except Exception as e:
            logger.info(f"Error in {filename}: {e}")
            logger.info(traceback.format_exc())
            os.remove(temp_path)
     

    # chunking recursively without semantic search, this does not uses openai embeddings for chunking
    text_splitter = RecursiveCharacterTextSplitter(
    chunk_size = 1536,
    chunk_overlap  = 0,
    length_function = len
    )
    # uses openai embeddings for chunking, costs time and money but gives good performances, not ideal for real-time
    # text_splitter = SemanticChunker(OpenAIEmbeddings(), breakpoint_threshold_type="percentile", number_of_chunks= 10000)
    logger.info("Before Embeddings!",)
    
    logger.info("Now Doing Embeddings!")
    
    docsearch = None

    try:
        if texts: # for pptx, docx, pdf
            logger.info("Running Text Merged (Pdfs other than extractedcontent.pdf)")
            text_splitter_formerged = RecursiveCharacterTextSplitter(chunk_size=1536, chunk_overlap=128, length_function=len)
            text_chunks = text_splitter_formerged.split_text(texts)

            batch_size = 100  # As per the limit mentioned in your error
            total_chunks = len(text_chunks)
            logger.info(f"total_chunks: {total_chunks}",)

            for i in range(0, total_chunks, batch_size):
                batch = text_chunks[i:i + batch_size]

                docsearch = FAISS.from_texts(batch, embeddings)
                logger.info(f"docsearch made for batch,{len(batch)} of the total {total_chunks}")

        elif raw_text:
            raw_text_splitted_chunks = text_splitter.split_text(raw_text)

            batch_size = 100  # As per the limit mentioned in your error
            total_chunks = len(raw_text_splitted_chunks)
            logger.info(f"total_chunks: {total_chunks}",)

            for i in range(0, total_chunks, batch_size):
                batch_raw_text = raw_text_splitted_chunks[i:i + batch_size]

                docsearch = FAISS.from_texts(batch_raw_text, embeddings)
                logger.info(f"docsearch made for batch,{len(batch_raw_text)} of the total {total_chunks}")

    except Exception as e:
        logger.error(f"Error with vectorization:{str(e)}")

    logger.info("docsearch made")
    return docsearch

def URL_IMG_EXTRACT(soup, session_var, base_url):
    output_path_byfile = f"./imagefolder_{session_var}/images_{session_var}_URL"
    if not os.path.exists(output_path_byfile):
        os.makedirs(output_path_byfile, exist_ok=True) 
    index = 1
    skip_patterns = ['/icons/', '/logos/', '/ads/', '/footer/']
    for image in soup.find_all('img'):
        try:
            logger.info(f"{image}")
            image_src = image['src']
            logger.info(f"image_src: {image_src}",)
            extension_src = image_src.split('.')[-1]
            logger.info(f"extension image_src: {extension_src}",)
            if extension_src not in ['png', 'jpg', 'jpeg', 'JPG', 'webp']:
                continue
            if any(pattern in image_src for pattern in skip_patterns):
              continue
            if any('logo' in cls for cls in image.get('class', [])):
                continue
            full_image_url = urljoin(base_url, image_src)
            logger.info(f"Selected: {full_image_url}")
            img_filename = os.path.join(output_path_byfile, f'FileName URL ImageNumber {index}.{extension_src}')
            # if extension_src == 'svg':
            #     svg_filename = os.path.join(output_path_byfile, f'URL image_{index}.png')
            #     png = cairosvg.svg2png(url = full_image_url)
            #     with open(svg_filename, 'wb') as f:
            #         f.write(png)
            #         index += 1
            #         continue

            # if pil would want to be used for saving and converting
            # image = requests.get(full_image_url)
            # image_content = image.content
            # try:
            #     img = Image.open(io.BytesIO(image_content))
            #     img.save(img_filename)
            #     index += 1
            # except:
            #     print(f"Failed to download {full_image_url}")
            
            with open(img_filename, 'wb') as f:
                image = urllib.request.urlopen(full_image_url).read()
                f.write(image)
                index += 1
        except Exception as e:
          logger.error(f"Failed to download {full_image_url}, error: {str(e)}")

def REMOVE_DUP_IMG(image_dir):
    # Directory containing images
    image_dir = image_dir

    # store hashes and their corresponding filenames
    hashes = {}
    duplicates = []

    # loop over the image files in the directory
    for root, dirs, files in os.walk(image_dir):
      for filename in files:
          if filename.endswith(('.png', '.jpg', '.jpeg', '.webp', '.JPG')):
            file_path = os.path.join(root, filename) # root returns the file path including subfolder
            logger.info(file_path)

            with Image.open(file_path) as img:
                width, height = img.size
                        
            if width * height <= 10000:
                os.remove(file_path)
                continue
            else:
                with Image.open(file_path) as img:

                    # generate the perceptual hash for the image
                    hash = imagehash.phash(img)

                    # seeing if this hash already exists in the dictionary
                    if hash in hashes:
                        duplicates.append(filename)
                        logger.info(f"Duplicate found: {filename} is a duplicate of {hashes[hash]}")
                        os.remove(file_path)  # Delete the duplicate file
                    else:
                        hashes[hash] = filename

    logger.info(f"Duplicate removal complete. Removed files: {duplicates}",)

def PRODUCE_LEARNING_OBJ_COURSE(query, docsearch, llm, model_type, language):
    logger.info("PRODUCE_LEARNING_OBJ_COURSE Initiated!")
    docs = docsearch.similarity_search(query, k=2)
    docs_main = " ".join([d.page_content for d in docs])
    logger.info(f"1st Docs_main of /Decide route:{docs_main}")

    if model_type=="gemini":
        chain = PROMPTS.prompt_LO_CA_GEMINI | llm.bind(generation_config={"response_mime_type": "application/json"})
        # chain = LLMChain(prompt=PROMPTS.prompt_LO_CA_GEMINI, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
    else:
        chain = PROMPTS.prompt_LO_CA_GEMINI | llm.bind(response_format={"type": "json_object"})
        # chain = LLMChain(prompt=PROMPTS.prompt_LO_CA_GEMINI, llm=llm.bind(response_format={"type": "json_object"}))
    
    logger.info("response_LO_CA started")
    response_LO_CA = chain.invoke({"input_documents": docs_main,"human_input": query, "language":language})
    logger.info(f"{response_LO_CA}")
    logger.info("response_LO_CA ended")

    return response_LO_CA

def PRODUCE_LEARNING_OBJ_COURSE_WITHOUT_FILE(query, llm, model_type, language):
    logger.info("PRODUCE_LEARNING_OBJ_COURSE_WITHOUT_FILE Initiated!")

    if model_type=="gemini":
        chain = PROMPTS_WITHOUT_FILE.prompt_LO_CA_GEMINI | llm.bind(generation_config={"response_mime_type": "application/json"})
        # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_LO_CA_GEMINI, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
    else:
        chain = PROMPTS_WITHOUT_FILE.prompt_LO_CA_GEMINI | llm.bind(response_format={"type": "json_object"})
        # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_LO_CA_GEMINI, llm=llm.bind(response_format={"type": "json_object"}))

    response_LO_CA = chain.invoke({"human_input": query, "language":language})
    
    return response_LO_CA

def RE_SIMILARITY_SEARCH(query, docsearch, output_path, model_type,model_name, summarize_images, language, llm_img_summary):
    logger.info("RE_SIMILARITY_SEARCH Initiated!")
    docs = docsearch.similarity_search(query, k=2)
    logger.info(f"docs from RE_SIMILARITY_SEARCH:\n{docs}",)
    if summarize_images == "on":
        logger.info(f"Tells me to summarize images, {summarize_images}")
        PageNumberList = []
        for relevant_doc in docs:
            relevant_doc = relevant_doc.page_content
            logger.info(relevant_doc)

            pattern_this_pptx = r"SlideNumber (\d+) FileName (.+?) ---"
            # Find all matches for "[This Page is PageNumber:]"
            matches_this_pptx = re.findall(pattern_this_pptx, relevant_doc)

            pattern_this_end_pptx = r"End of SlideNumber (\d+) with Filename (.+?) ---"
            # Find all matches for "[This Page is PageNumber:]"
            matches_this_end_pptx = re.findall(pattern_this_end_pptx, relevant_doc)


            pattern_this_doc = r'----media/ImageNumber (\d+) PageNumber Null of FileName (.+)----'
            matches_this_doc = re.findall(pattern_this_doc, relevant_doc)
            
            pattern_end = r'End of PageNumber (\d+) of file name (.+)\n'
            pattern_this_page = r'The Content of PageNumber (\d+) of file name (.+) is:\n'
            # Find all matches for "End of PageNumber:"
            matches_end = re.findall(pattern_end, relevant_doc)

            # Find all matches for "[This Page is PageNumber:]"
            matches_this_page = re.findall(pattern_this_page, relevant_doc)

            # Combine the matches
            for num in matches_this_page + matches_end + matches_this_doc + matches_this_pptx + matches_this_end_pptx:
                PageNumberList.append(num)

            PageNumberList = list(set(PageNumberList))
            logger.info(f"PageNumberList:\n{PageNumberList}",)

        image_elements = []
        image_summaries = []
        pixmapped_image_summaries = []

        def encode_image(image_path):
            basename = os.path.basename(image_path)
            with Image.open(image_path) as img:
                width, height = img.size
                logger.info(f"{basename} size is {width},{height}")

                if width*height <= 10000:
                    logger.info(f"Skipping {basename} as it is smaller than 100x100 pixels with size as {width},{height}")
                    return None

                if width*height > 262144:
                    # Resize the image
                    img = img.resize((512, 512))
                    # Save the resized image to a temporary file
                    basenama = os.path.basename(image_path)
                    extensiona = basenama.rsplit('.', 1)[1].lower()

                    temp_patha = image_path + f"_temp_img.{extensiona}"
                    img.save(temp_patha)
                
                    # Encode the resized image
                    with open(temp_patha, "rb") as f:
                        encoded_image = base64.b64encode(f.read()).decode('utf-8')

                    # Remove the temporary file
                    os.remove(temp_patha)
                else:
                    logger.info(f"{basename} is less than 262144 having {width}, {height}")
                    with open(image_path, "rb") as f:
                        encoded_image = base64.b64encode(f.read()).decode('utf-8')

                return encoded_image


        def summarize_image(encoded_image, basename, language):
            logger.info(f"language is: {language}")
            prompt = [
                SystemMessage(content="You are a bot that is good at analyzing images."),
                HumanMessage(content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.", },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{encoded_image}"
                        },
                    },
                ])
            ]

            prompt_gemini = HumanMessage(
                content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.", },
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{encoded_image}"},
                ]
            )
            
            if model_type == 'gemini':
                logger.info("Gemini summarizing images NOW")
                response = ChatGoogleGenerativeAI(model=model_name,temperature=0,max_output_tokens=250).invoke([prompt_gemini])
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
                
            else:
                logger.info("Openai summarizing images NOW")
                response = AzureChatOpenAI(deployment_name=model_name, temperature=0, max_tokens=250,
                                            openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")).invoke(prompt)
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
            
            chain = PROMPTS.prompt_polish_summary | llm_img_summary
            # chain = LLMChain(prompt=PROMPTS.prompt_polish_summary,llm=llm_img_summary)
            polish_summary = chain.invoke({"basename": basename,"description": response.content,"language":language})
            
            return polish_summary.content
            

        def url_summarize_image(encoded_image, basename, language):
            prompt = [
                SystemMessage(content="You are a bot that is good at analyzing images."),
                HumanMessage(content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.",
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{encoded_image}"
                        },
                    },
                ])
            ]

            prompt_gemini = HumanMessage(
                content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.",
                    },  # You can optionally provide text parts
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{encoded_image}"},
                ]
            )

            if model_type == 'gemini':
                logger.info("Gemini summarizing images NOW")
                response = ChatGoogleGenerativeAI(model=model_name,temperature=0,max_output_tokens=250).invoke([prompt_gemini])
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
                
            else:
                logger.info("Openai summarizing images NOW")
                response = AzureChatOpenAI(deployment_name=model_name, temperature=0, max_tokens=250,
                                            openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")).invoke(prompt)
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
            
            chain = PROMPTS.prompt_polish_summary | llm_img_summary
            # chain = LLMChain(prompt=PROMPTS.prompt_polish_summary,llm=llm_img_summary)
            polish_summary = chain.invoke({"basename": basename,"description": response.content,"language":language})
            
            return polish_summary.content

        def pixmapped_summarize_image(encoded_image, basename, language):
            
            prompt = [
                SystemMessage(content="You are a bot that describes only the figures present in this pdf page's image."),
                HumanMessage(content=[
                    {
                        "type": "text",
                        "text": f"Describe in a brief manner only the figures present in this pdf page's image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.", },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{encoded_image}"
                        },
                    },
                ])
            ]

            prompt_gemini = HumanMessage(
                content=[
                    {
                        "type": "text",
                        "text": f"Describe in a brief manner only the figures present in this pdf page's image in the language of {language}, since your responses are given to {language} speakers and they can only understand the language of {language}.", },
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{encoded_image}"},
                ]
            )
            
            if model_type == 'gemini':
                logger.info("Gemini summarizing images NOW")
                response = ChatGoogleGenerativeAI(model=model_name,temperature=0,max_output_tokens=2048).invoke([prompt_gemini])
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
                
            else:
                logger.info("Openai summarizing images NOW")
                response = AzureChatOpenAI(deployment_name=model_name, temperature=0, max_tokens=2048,
                                            openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")).invoke(prompt)
                img_desc = response.content
                logger.info(f"Img Summary is: {type(img_desc)}/n{img_desc}")
            
            return img_desc

        for root, dirs, files in os.walk(output_path):
            for i in files:
                if i.endswith(('.png', '.jpg', '.jpeg', '.webp', '.JPG')):

                    if "URL" in i:
                        image_path = os.path.join(root, i)
                        basename = os.path.basename(image_path)
                        logger.info(os.path.basename(image_path))
                        encoded_image = encode_image(image_path)
                        if encoded_image is not None:
                            image_elements.append(encoded_image)
                            summary = url_summarize_image(encoded_image,basename, language)
                            image_summaries.append(summary)

                    else:
                        for page_number, file in PageNumberList:
                            if f"FileName {file} PageNumber {page_number}" in i:
                                image_path = os.path.join(root, i)
                                basename = os.path.basename(image_path)
                                logger.info(os.path.basename(image_path))
                                if "Pixmapped" in basename:
                                    logger.info(f"Pixmap detected for : {basename}")
                                    encoded_image = encode_image(image_path)
                                    if encoded_image is not None:
                                        image_elements.append(encoded_image)
                                        summary = pixmapped_summarize_image(encoded_image,basename, language)
                                        pixmapped_image_summaries.append(summary)
                                        image_summaries_string = "\n".join(pixmapped_image_summaries) #convert list to string to add in the langchain Document data type                                                                        
                                        docs.append(Document(page_content=f"\nThe Content of PageNumber {page_number} of file name {file} is:\nThis PIXMAPPED IMAGE described as: {image_summaries_string}.\nEnd of PageNumber {page_number} of file name {file}\n"))                                        
                                        logger.info(f"The docs are:\n{docs}")
                                    # os.remove(image_path) # WARNING: THIS MIGHT EFFECT FIND_IMAGES!!!!!!

                                else:
                                    encoded_image = encode_image(image_path)
                                    if encoded_image is not None:
                                        image_elements.append(encoded_image)
                                        summary = summarize_image(encoded_image,basename, language)
                                        image_summaries.append(summary)
                            elif f"FileName {file} PageNumber Null ImageNumber {page_number}" in i:
                                image_path = os.path.join(root, i)
                                basename = os.path.basename(image_path)
                                logger.info(os.path.basename(image_path))
                                encoded_image = encode_image(image_path)
                                if encoded_image is not None:
                                    image_elements.append(encoded_image)
                                    summary = summarize_image(encoded_image,basename, language)
                                    image_summaries.append(summary)
                            elif f"FileName {file} SlideNumber {page_number}" in i:
                                image_path = os.path.join(root, i)
                                basename = os.path.basename(image_path)
                                logger.info(os.path.basename(image_path))
                                encoded_image = encode_image(image_path)
                                if encoded_image is not None:
                                    image_elements.append(encoded_image)
                                    summary = summarize_image(encoded_image,basename, language)
                                    image_summaries.append(summary)

        logger.info(f"image_summaries::\n{image_summaries}",)

        image_summaries_string = "\n".join(image_summaries) #convert list to string to add in the langchain Document data type
        docs.append(Document(page_content=f"Useful Image/s for all the above content::\n{image_summaries_string}"))

    return docs

mpv_list = ["NO number of MediaBlock/s and ONLY TextBlock/s", "more TextBlock/s compared to MediaBlock/s",
       "BALANCED number of MediaBlock/s compared to TextBlock/s",
       "more MediaBlock/s compared to TextBlock/s", "ONLY MediaBlock/s and NO number of TextBlock/s"]

def TALK_WITH_RAG(scenario, content_areas, learning_obj, query, docs_main, llm, model_type, model_name,embeddings, language, mpv):
    logger.info("TALK_WITH_RAG Initiated!")
    # if we are getting docs_main already from the process_data flask route then comment, else
    # UNcomment if you want more similarity_searching based on Learning obj and content areas!
    # docs = docsearch.similarity_search(query, k=3)
    # docs_main = " ".join([d.page_content for d in docs])
    
    mpv_string = mpv_list[int(mpv)]
    logger.info(f"mpv list string is: {mpv_string}")

    responses = ''
    def is_json_parseable(json_string):
        try:
            json_object = json.loads(json_string)
        except ValueError as e:
            return False, str(e)
        return True, json_object

    if scenario == "auto":
        logger.info(f"SCENARIO ====PROMPT{scenario}",)
        # chain = prompt | llm | {f"{llm_memory}": RunnablePassthrough()}
        

        ### SEMANTIC ROUTES LOGIC ###
        if model_type == 'gemini':
            llm_auto = ChatGoogleGenerativeAI(model=model_name,temperature=0.4, max_output_tokens=32) 
            # llm_auto_chain = LLMChain(prompt=PROMPTS.promptSelector, llm=llm_auto.bind(generation_config={"response_mime_type": "application/json"})) 
            llm_auto_chain = PROMPTS.promptSelector | llm_auto.bind(generation_config={"response_mime_type": "application/json"})
        else:
            llm_auto =  AzureChatOpenAI(deployment_name=model_name, temperature=0.4, max_tokens=32,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )
            # llm_auto_chain = LLMChain(prompt=PROMPTS.promptSelector, llm=llm_auto.bind(response_format={"type": "json_object"}))
            llm_auto_chain = PROMPTS.promptSelector | llm_auto.bind(response_format={"type": "json_object"})

        selected = llm_auto_chain.invoke({"input_documents": docs_main, "human_input": query})

        logger.info(f"Semantic Scenario Selected of NAME: {selected.content}",)

        selected = json.loads(selected.content)
        max_similarity = selected["Bot"]
        logger.info(f"max_similarity is:{max_similarity}")

        ############################

        if max_similarity == "Gamified Scenario":
            logger.info("Gamified Auto Selected")
            scenario = "gamified"

        elif max_similarity == "Linear Scenario":
            logger.info("Linear Auto Selected")
            scenario = "linear"

        elif max_similarity == "Simulation Scenario":
            logger.info(f"Simulation Auto Selected")
            scenario = "simulation"

        elif max_similarity == "Branched Scenario":
            logger.info(f"Branched Auto Selected")
            scenario = "branched"

        else:
            logger.info(f"AUTO SELECTION FAILED, Selecting Default Scenario of LINEAR SCENARIO")
            scenario = "linear"


    if scenario == "linear":
        logger.info(f"SCENARIO ====prompt_linear : {scenario}")
        if model_type == 'gemini':
            chain = PROMPTS.prompt_linear | llm.bind(generation_config={"response_mime_type": "application/json"})
            # chain = LLMChain(prompt=PROMPTS.prompt_linear,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
        else:
            chain = PROMPTS.prompt_linear | llm
            # chain = LLMChain(prompt=PROMPTS.prompt_linear,llm=llm)   

        response = chain.invoke({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}")
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}")
            logger.info(f"changed:::\n{modified_txt}")

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}")
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}")

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                chain_retry = PROMPTS.prompt_linear_retry | llm.bind(generation_config={"response_mime_type": "application/json"})
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_linear_retry,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            else:
                chain_retry = PROMPTS.prompt_linear_retry | llm
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_linear_retry,llm=llm)

            response_retry = chain_retry.invoke({"incomplete_response": responses, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    chain_simplify = PROMPTS.prompt_linear_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_linear_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                else:
                    chain_simplify = PROMPTS.prompt_linear_simplify | llm
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_linear_simplify,llm=llm)

                response_retry_simplify = chain_simplify.invoke({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    

    elif scenario == "branched":
        logger.info(f"SCENARIO ====branched : {scenario}",)
        
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
            llm_setup_continue = ChatGoogleGenerativeAI(model=model_name,temperature=0.1)
            # chain = LLMChain(prompt=PROMPTS.prompt_branched,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS.prompt_branched | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            llm_setup = AzureChatOpenAI(deployment_name=model_name, temperature=0,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )  
            llm_setup_continue = AzureChatOpenAI(deployment_name=model_name, temperature=0.1,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )
            # chain = LLMChain(prompt=PROMPTS.prompt_branched,llm=llm)  
            chain = PROMPTS.prompt_branched | llm 

        if model_type == 'gemini':
            chain1 = PROMPTS.prompt_branched_setup | llm_setup
            # chain1 = LLMChain(prompt=PROMPTS.prompt_branched_setup,llm=llm_setup)
        else:
            chain1 = PROMPTS.prompt_branched_setup | llm_setup
            # chain1 = LLMChain(prompt=PROMPTS.prompt_branched_setup,llm=llm_setup)

        response1 = chain1.invoke({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
        if "[END_OF_RESPONSE]" not in response1.content:
            count_setup_retry = 0
            while "[END_OF_RESPONSE]" not in response1.content and count_setup_retry<=3:
                logger.info("[END_OF_RESPONSE] not found")
                contd_response1 = response1.content + "[CONTINUE_EXACTLY_FROM_HERE]"
                chain_setup_retry = PROMPTS.prompt_branched_setup_continue | llm_setup_continue
                # chain_setup_retry = LLMChain(prompt=PROMPTS.prompt_branched_setup_continue,llm=llm_setup_continue)
                response1 = chain_setup_retry.invoke({"past_response": contd_response1,"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
                logger.info(f"CONTINUED Response 1 IS::\n{response1.content}")
                response1.content = contd_response1 + response1.content
                response1.content = re.sub(r'\[CONTINUE_EXACTLY_FROM_HERE\]', ' ', response1.content)
                
                logger.info(f"JOINED Response 1 IS::\n{response1.content}")
                count_setup_retry += 1
        else:
            logger.info(f"Response 1 is::\n{response1.content}",)


        response = chain.invoke({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        # clean_json = response['text'].strip('`json ')
        # response = {'text':clean_json}
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}", )
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}",)
            logger.info(f"changed:::\n{modified_txt}",)

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}", )

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                chain_retry = PROMPTS.prompt_branched_retry | llm.bind(generation_config={"response_mime_type": "application/json"})
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_branched_retry,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            else:
                chain_retry = PROMPTS.prompt_branched_retry | llm
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_branched_retry,llm=llm)

            response_retry = chain_retry.invoke({"incomplete_response": responses,"micro_subtopics":response1.content, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!:\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    chain_simplify = PROMPTS.prompt_branched_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_branched_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                else:
                    chain_simplify = PROMPTS.prompt_branched_simplify | llm
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_branched_simplify,llm=llm)

                response_retry_simplify = chain_simplify.invoke({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    

    elif scenario == "simulation":
        logger.info(f"SCENARIO ====prompt_simulation_pedagogy : {scenario}",)
        # summarized first, then response
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0.3)
            llm_setup_continue = ChatGoogleGenerativeAI(model=model_name,temperature=0.1)
            # chain = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS.prompt_simulation_pedagogy_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            llm_setup = AzureChatOpenAI(deployment_name=model_name, temperature=0.3,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )  
            llm_setup_continue = AzureChatOpenAI(deployment_name=model_name, temperature=0.1,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        ) 
            # chain = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_gemini,llm=llm) 
            chain = PROMPTS.prompt_simulation_pedagogy_gemini | llm  

        if model_type == 'gemini':
            chain1 = PROMPTS.prompt_simulation_pedagogy_setup | llm_setup
            # chain1 = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_setup,llm=llm_setup)
        else:
            chain1 = PROMPTS.prompt_simulation_pedagogy_setup | llm_setup
            # chain1 = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_setup,llm=llm_setup)

        response1 = chain1.invoke({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
        if "[END_OF_RESPONSE]" not in response1.content:
            count_setup_retry = 0
            while "[END_OF_RESPONSE]" not in response1.content and count_setup_retry<=3:
                logger.info("[END_OF_RESPONSE] not found")
                contd_response1 = response1.content + "[CONTINUE_EXACTLY_FROM_HERE]"
                chain_setup_retry = PROMPTS.prompt_simulation_pedagogy_setup_continue | llm_setup_continue
                # chain_setup_retry = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_setup_continue,llm=llm_setup_continue)
                response1 = chain_setup_retry.invoke({"past_response": contd_response1,"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
                logger.info(f"CONTINUED Response 1 IS::\n{response1.content}")
                response1.content = contd_response1 + response1.content
                response1.content = re.sub(r'\[CONTINUE_EXACTLY_FROM_HERE\]', ' ', response1.content)
                
                logger.info(f"JOINED Response 1 IS::\n{response1.content}")
                count_setup_retry += 1
        else:
            logger.info(f"Response 1 is::\n{response1.content}",)


        response = chain.invoke({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}", )
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}",)
            logger.info(f"changed:::\n{modified_txt}",)

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}", )

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                chain_retry = PROMPTS.prompt_simulation_pedagogy_retry_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_retry_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            else:
                chain_retry = PROMPTS.prompt_simulation_pedagogy_retry_gemini | llm
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_retry_gemini,llm=llm)

            response_retry = chain_retry.invoke({"incomplete_response": responses,"simulation_story":response1.content, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!:\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    chain_simplify = PROMPTS.prompt_simulation_pedagogy_gemini_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_gemini_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                else:
                    chain_simplify = PROMPTS.prompt_simulation_pedagogy_gemini_simplify | llm
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_simulation_pedagogy_gemini_simplify,llm=llm)

                response_retry_simplify = chain_simplify.invoke({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                            

    elif scenario == "gamified":
        logger.info(f"SCENARIO ====prompt_gamified : {scenario}",)
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
            llm_setup_continue = ChatGoogleGenerativeAI(model=model_name,temperature=0.1)
            # chain = LLMChain(prompt=PROMPTS.prompt_gamified_json,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS.prompt_gamified_json | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            llm_setup = AzureChatOpenAI(deployment_name=model_name, temperature=0,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )  
            llm_setup_continue = AzureChatOpenAI(deployment_name=model_name, temperature=0.1,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )
            # chain = LLMChain(prompt=PROMPTS.prompt_gamified_json,llm=llm)
            chain = PROMPTS.prompt_gamified_json | llm   

        if model_type == 'gemini':
            chain1 = PROMPTS.prompt_gamified_setup | llm_setup 
            # chain1 = LLMChain(prompt=PROMPTS.prompt_gamified_setup,llm=llm_setup)
        else:
            chain1 = PROMPTS.prompt_gamified_setup | llm_setup 
            # chain1 = LLMChain(prompt=PROMPTS.prompt_gamified_setup,llm=llm_setup)

        response1 = chain1.invoke({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
        if "[END_OF_RESPONSE]" not in response1.content:
            count_setup_retry = 0
            while "[END_OF_RESPONSE]" not in response1.content and count_setup_retry<=3:
                logger.info("[END_OF_RESPONSE] not found")
                contd_response1 = response1.content + "[CONTINUE_EXACTLY_FROM_HERE]"
                chain_setup_retry = PROMPTS.prompt_gamified_setup_continue | llm_setup_continue
                # chain_setup_retry = LLMChain(prompt=PROMPTS.prompt_gamified_setup_continue,llm=llm_setup_continue)
                response1 = chain_setup_retry.invoke({"past_response": contd_response1,"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language})
                logger.info(f"CONTINUED Response 1 IS::\n{response1.content}")
                response1.content = contd_response1 + response1.content
                response1.content = re.sub(r'\[CONTINUE_EXACTLY_FROM_HERE\]', ' ', response1.content)
                
                logger.info(f"JOINED Response 1 IS::\n{response1.content}")
                count_setup_retry += 1
        else:
            logger.info(f"Response 1 is::\n{response1.content}",)


        response = chain({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}", )
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}",)
            logger.info(f"changed:::\n{modified_txt}",)

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}", )

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                chain_retry = PROMPTS.prompt_gamified_pedagogy_retry_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_gamified_pedagogy_retry_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            else:
                chain_retry = PROMPTS.prompt_gamified_pedagogy_retry_gemini | llm
                # chain_retry = LLMChain(prompt=PROMPTS.prompt_gamified_pedagogy_retry_gemini,llm=llm)

            response_retry = chain_retry.invoke({"incomplete_response": responses,"exit_game_story":response1.content, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!:\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    chain_simplify = PROMPTS.prompt_gamify_pedagogy_gemini_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_gamify_pedagogy_gemini_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                else:
                    chain_simplify = PROMPTS.prompt_gamify_pedagogy_gemini_simplify | llm
                    # chain_simplify = LLMChain(prompt=PROMPTS.prompt_gamify_pedagogy_gemini_simplify,llm=llm)

                response_retry_simplify = chain_simplify.invoke({"response_of_bot": response1.content,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    
     
    logger.info(f"The output is as follows::\n{response.content}",)
    return response.content, scenario

def TALK_WITH_RAG_WITHOUT_FILE(scenario, content_areas, learning_obj, query, llm, model_type, model_name, language, mpv):
    mpv_string = mpv_list[int(mpv)]
    logger.info(f"mpv list string is: {mpv_string}")

    responses = ''
    def is_json_parseable(json_string):
        try:
            json_object = json.loads(json_string)
        except ValueError as e:
            return False, str(e)
        return True, json_object

    if scenario == "auto":
        logger.info(f"SCENARIO ====PROMPT{scenario}",)
        # chain = prompt | llm | {f"{llm_memory}": RunnablePassthrough()}
        

        ### SEMANTIC ROUTES LOGIC ###
        if model_type == 'gemini':
            llm_auto = ChatGoogleGenerativeAI(model=model_name,temperature=0.4, max_output_tokens=32) 
            # llm_auto_chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.promptSelector, llm=llm_auto.bind(generation_config={"response_mime_type": "application/json"})) 
            llm_auto_chain = PROMPTS_WITHOUT_FILE.promptSelector | llm_auto.bind(generation_config={"response_mime_type": "application/json"})
        else:
            llm_auto =  AzureChatOpenAI(deployment_name=model_name, temperature=0.4, max_tokens=32,
                                        openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                        )
            # llm_auto_chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.promptSelector, llm=llm_auto.bind(response_format={"type": "json_object"}))
            llm_auto_chain = PROMPTS_WITHOUT_FILE.promptSelector | llm_auto.bind(response_format={"type": "json_object"})

        selected = llm_auto_chain.invoke({"human_input": query, "content_areas": content_areas, "learning_objectives": learning_obj})

        logger.info(f"Semantic Scenario Selected of NAME: {selected.content}",)

        selected = json.loads(selected.content)
        max_similarity = selected["Bot"]
        logger.info(f"max_similarity is:{max_similarity}")

        ############################

        if max_similarity == "Gamified Scenario":
            logger.info("Gamified Auto Selected")
            scenario = "gamified"

        elif max_similarity == "Linear Scenario":
            logger.info("Linear Auto Selected")
            scenario = "linear"

        elif max_similarity == "Simulation Scenario":
            logger.info(f"Simulation Auto Selected")
            scenario = "simulation"

        elif max_similarity == "Branched Scenario":
            logger.info(f"Branched Auto Selected")
            scenario = "branched"

        else:
            logger.info(f"AUTO SELECTION FAILED, Selecting Default Scenario of LINEAR SCENARIO")
            scenario = "linear"


    if scenario == "linear":
        logger.info(f"SCENARIO ====prompt_linear : {scenario}")
        if model_type == 'gemini':
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS_WITHOUT_FILE.prompt_linear | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear,llm=llm)   
            chain = PROMPTS_WITHOUT_FILE.prompt_linear | llm

        response = chain.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}")
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}")
            logger.info(f"changed:::\n{modified_txt}")

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}")
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}")

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_retry,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_linear_retry | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_retry,llm=llm)
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_linear_retry | llm

            response_retry = chain_retry.invoke({"incomplete_response": responses, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_linear_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_simplify,llm=llm)
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_linear_simplify | llm

                response_retry_simplify = chain_simplify.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    

    elif scenario == "branched":
        logger.info(f"SCENARIO ====branched : {scenario}",)
        
        if model_type == 'gemini':
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS_WITHOUT_FILE.prompt_branched | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched,llm=llm)   
            chain = PROMPTS_WITHOUT_FILE.prompt_branched | llm

        response = chain.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}")
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}")
            logger.info(f"changed:::\n{modified_txt}")

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}")
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}")

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_retry,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_branched_retry | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_retry,llm=llm)
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_branched_retry | llm

            response_retry = chain_retry.invoke({"incomplete_response": responses, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_branched_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_simplify,llm=llm)
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_branched_simplify | llm

                response_retry_simplify = chain_simplify.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    

    elif scenario == "simulation":
        logger.info(f"SCENARIO ====prompt_simulation_pedagogy : {scenario}",)
        if model_type == 'gemini':
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini,llm=llm)   
            chain = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini | llm

        response = chain.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}")
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}")
            logger.info(f"changed:::\n{modified_txt}")

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}")
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}")

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_retry_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_retry_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_retry_gemini,llm=llm)
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_retry_gemini | llm

            response_retry = chain_retry.invoke({"incomplete_response": responses, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini_simplify,llm=llm)
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_simulation_pedagogy_gemini_simplify | llm
                
                response_retry_simplify = chain_simplify.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                            

    elif scenario == "gamified":
        logger.info(f"SCENARIO ====prompt_gamified : {scenario}",)
        if model_type == 'gemini':
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamified_json,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
            chain = PROMPTS_WITHOUT_FILE.prompt_gamified_json | llm.bind(generation_config={"response_mime_type": "application/json"})
        else:
            # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamified_json,llm=llm)   
            chain = PROMPTS_WITHOUT_FILE.prompt_gamified_json | llm

        response = chain.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
        
        is_valid, result = is_json_parseable(response.content)
        countd=1
        while not is_valid and countd<=2:
            txt = response.content
            logger.info(f"CHAIN_RETRY BEGINS for the failed response:\n{txt}")
            ### REGEX to remove last incomplete id block ###
            modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
            if modified_txt:
                modified_txt = modified_txt[0]  # Get the matched string
            else:
                modified_txt = txt  # No match found, return original
            logger.info(f"original:::\n{txt}")
            logger.info(f"changed:::\n{modified_txt}")

            # Finding if corrupt edges exists further and to remove it via if loop
            find_edges = re.findall(r'.*}, "edges": \[', modified_txt, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}")
                # Using regex to replace the specific pattern
                modified_txt = re.sub(r'}(?=, "edges": \[)', '}]', modified_txt, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{modified_txt}")

            responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
            logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

            if model_type == 'gemini':
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamified_pedagogy_retry_gemini,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_gamified_pedagogy_retry_gemini | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                # chain_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamified_pedagogy_retry_gemini,llm=llm)
                chain_retry = PROMPTS_WITHOUT_FILE.prompt_gamified_pedagogy_retry_gemini | llm

            response_retry = chain_retry.invoke({"incomplete_response": responses, "language":language, "mpv":mpv, "mpv_string":mpv_string})
            logger.info(f"response contd... is:\n{response_retry.content}",)

            responses = modified_txt + response_retry.content #changed modified_text to responses
            logger.info(f"responses+continued Combined is:\n{responses}",)

            # Finding if corrupt edges exists AFTER combined prompts
            find_edges = re.findall(r'.*}, "edges": \[', responses, re.DOTALL)
            if find_edges:
                find_edges = find_edges[0]  # Get the matched string
                logger.info(f"Corrupt edges found:\n{find_edges}",)
                # Using regex to replace the specific pattern
                responses = re.sub(r'}(?=, "edges": \[)', '}]', responses, flags=re.DOTALL)
                logger.info(f"Corrected corrupt edges:\n{responses}", )

            response.content = responses

            is_valid, result = is_json_parseable(responses)
            logger.info(f"Parseability status:\n{result}", )
            countd+=1
            logger.info(f"contd count is:\n{countd}",)

        if is_valid == False and countd==3: #countd==4 shows while loop has exited with failure 
            logger.info(f"The retry is also not parseable!\n{responses}", )
            max_attempts = 1  # Maximum number of attempts
            attempts = 1
            while attempts <= max_attempts:

                if model_type == 'gemini':
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamify_pedagogy_gemini_simplify,llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_gamify_pedagogy_gemini_simplify | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain_simplify = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamify_pedagogy_gemini_simplify,llm=llm)
                    chain_simplify = PROMPTS_WITHOUT_FILE.prompt_gamify_pedagogy_gemini_simplify | llm

                response_retry_simplify = chain_simplify.invoke({"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify.content)
                if is_valid_retry_simplify == True:
                    response.content = response_retry_simplify.content
                    logger.info(f"Result successfull for simplified response:\n{response.content}",)
                    break
                else:
                    logger.info(f"Attempt {attempts} also failed to parse JSON. Error:\n {response_retry_simplify.content}")
                    attempts += 1
                    
     
    logger.info(f"The output is as follows::\n{response.content}",)
    return response.content


def REPAIR_SHADOW_EDGES(scenario, original_txt,model_type, model_name, language, mpv, repair_shadows_without_file='0'):
    txt_output = None

    mpv_string = mpv_list[int(mpv)]
    logger.info(f"mpv list string is: {mpv_string}")

    if model_type == 'gemini':
        llm = ChatGoogleGenerativeAI(model=model_name,temperature=0)
    else:
        llm = AzureChatOpenAI(deployment_name=model_name, temperature=0,
                                    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION")
                                )  
        
    def is_json_parseable(json_string):
        try:
            json_object = json.loads(json_string)
        except ValueError as e:
            return False, str(e)
        return True, json_object

    def validate_edges(json_data):
        # txt_string = json.dumps(json_data)  # Convert dictionary to JSON string
        # json_data = json.loads(txt_string)  # Convert JSON string into dictionary
        json_data = json.loads(json_data)  # Convert JSON string into dictionary
        node_ids = {node['id'] for node in json_data['nodes']}
        error_flag = False

        for edge in json_data['edges']:
            source_exists = edge['source'] in node_ids
            target_exists = edge['target'] in node_ids

            if not source_exists or not target_exists:
                logger.info(f"Error occured:\n{edge}")
                edge['SHADOW EDGE BLOCK'] = 'SHADOW EDGES IN THIS BLOCK'  # Add error directly to the edge
                error_flag = True
                # If you want to find all errors, remove the break statement
                # break

        shadow_result = json.dumps(json_data, indent=4)
        return shadow_result, error_flag

    output, error_flag = validate_edges(original_txt)
    logger.info(f"error_flag: {error_flag}")

    if error_flag == True:

        logger.info(f"Error flag is: {error_flag} and so output is:\n{output}")

        if scenario == "linear":

            if model_type=="gemini":
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_linear_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain = PROMPTS.prompt_linear_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
                    chain = PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_linear_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS.prompt_linear_shadow_edges | llm.bind(response_format={"type": "json_object"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges | llm.bind(response_format={"type": "json_object"})

            shadow_response = chain.invoke({"output": output,"language":language, "mpv":mpv, "mpv_string":mpv_string})
            is_valid, result = is_json_parseable(shadow_response.content)
            countd=0
            while not is_valid and countd<=3:
                txt = shadow_response.content
                modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
                if modified_txt:
                    modified_txt = modified_txt[0]  # Get the matched string
                else:
                    modified_txt = txt  # No match found, return original
                logger.info(f"original:::\n{txt}")
                logger.info(f"changed:::\n{modified_txt}")

                responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
                logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

                if repair_shadows_without_file != "1":
                    # chain_edges_retry = LLMChain(prompt=PROMPTS.prompt_linear_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS.prompt_linear_shadow_edges_retry | llm
                else:
                    # chain_edges_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS_WITHOUT_FILE.prompt_linear_shadow_edges_retry | llm

                response_retry = chain_edges_retry.invoke({"incomplete_response": responses, "output":output, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                logger.info(f"response contd... is:\n{response_retry.content}",)

                responses = modified_txt + response_retry.content #changed modified_text to responses
                logger.info(f"responses+continued Combined is:\n{responses}",)
                
                shadow_response.content = responses

                is_valid, result = is_json_parseable(shadow_response.content)
                logger.info(f"Parseability status:\n{result}", )
                countd+=1
                logger.info(f"contd count is:\n{countd}",)

            logger.info("Success shadow repair!:",shadow_response.content)
            shadow_response = shadow_response.content
            logger.info(f"shadow_response type before: {type(shadow_response)}")
            logger.info(f"output type before: {type(output)}")

            shadow_response = json.loads(shadow_response)  # Convert JSON string into dictionary
            output = json.loads(output)  # Convert JSON string into dictionary

            logger.info(f"shadow_response type after: {type(shadow_response)} and output type after {type(output)}")

            output['edges']  = shadow_response['edges']  
            logger.info(f"{output} of {type(output)}")
            output = json.dumps(output, indent=2) # converts to dict to str
            is_valid_output, result = is_json_parseable(output)

            if is_valid_output == False:
                output = original_txt
                logger.info(f"The output was not parseable hence reverting to original_txt to this response:\n{output}")


        elif scenario == "branched":

            if model_type=="gemini":
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_branched_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
                    chain = PROMPTS.prompt_branched_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_branched_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS.prompt_branched_shadow_edges | llm.bind(response_format={"type": "json_object"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges | llm.bind(response_format={"type": "json_object"})

            shadow_response = chain.invoke({"output": output,"language":language, "mpv":mpv, "mpv_string":mpv_string})
            is_valid, result = is_json_parseable(shadow_response.content)
            countd=0
            while not is_valid and countd<=3:
                txt = shadow_response.content
                modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
                if modified_txt:
                    modified_txt = modified_txt[0]  # Get the matched string
                else:
                    modified_txt = txt  # No match found, return original
                logger.info(f"original:::\n{txt}")
                logger.info(f"changed:::\n{modified_txt}")

                responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
                logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

                if repair_shadows_without_file != "1":
                    # chain_edges_retry = LLMChain(prompt=PROMPTS.prompt_branched_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS.prompt_branched_shadow_edges_retry | llm
                else:
                    # chain_edges_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS_WITHOUT_FILE.prompt_branched_shadow_edges_retry | llm
                
                response_retry = chain_edges_retry.invoke({"incomplete_response": responses, "output":output, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                logger.info(f"response contd... is:\n{response_retry.content}",)

                responses = modified_txt + response_retry.content #changed modified_text to responses
                logger.info(f"responses+continued Combined is:\n{responses}",)
                
                shadow_response.content = responses

                is_valid, result = is_json_parseable(shadow_response.content)
                logger.info(f"Parseability status:\n{result}", )
                countd+=1
                logger.info(f"contd count is:\n{countd}",)

            logger.info("Success shadow repair!:",shadow_response.content)
            shadow_response = shadow_response.content
            logger.info(f"shadow_response type before: {type(shadow_response)}")
            logger.info(f"output type before: {type(output)}")

            shadow_response = json.loads(shadow_response)  # Convert JSON string into dictionary
            output = json.loads(output)  # Convert JSON string into dictionary

            logger.info(f"shadow_response type after: {type(shadow_response)} and output type after {type(output)}")

            output['edges']  = shadow_response['edges']  
            logger.info(f"{output} of {type(output)}")
            output = json.dumps(output, indent=2) # converts to dict to str
            is_valid_output, result = is_json_parseable(output)

            if is_valid_output == False:
                output = original_txt
                logger.info(f"The output was not parseable hence reverting to original_txt to this response:\n{output}")


        elif scenario == "simulation":

            if model_type=="gemini":
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_simulation_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
                    chain = PROMPTS.prompt_simulation_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_simulation_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS.prompt_simulation_shadow_edges | llm.bind(response_format={"type": "json_object"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges | llm.bind(response_format={"type": "json_object"})

            shadow_response = chain.invoke({"output": output,"language":language, "mpv":mpv, "mpv_string":mpv_string})
            is_valid, result = is_json_parseable(shadow_response.content)
            countd=0
            while not is_valid and countd<=3:
                txt = shadow_response.content
                modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
                if modified_txt:
                    modified_txt = modified_txt[0]  # Get the matched string
                else:
                    modified_txt = txt  # No match found, return original
                logger.info(f"original:::\n{txt}")
                logger.info(f"changed:::\n{modified_txt}")

                responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
                logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

                if repair_shadows_without_file != "1":
                    # chain_edges_retry = LLMChain(prompt=PROMPTS.prompt_simulation_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS.prompt_simulation_shadow_edges_retry | llm
                else:
                    # chain_edges_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS_WITHOUT_FILE.prompt_simulation_shadow_edges_retry | llm

                response_retry = chain_edges_retry.invoke({"incomplete_response": responses, "output":output, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                logger.info(f"response contd... is:\n{response_retry.content}",)

                responses = modified_txt + response_retry.content #changed modified_text to responses
                logger.info(f"responses+continued Combined is:\n{responses}",)
                
                shadow_response.content = responses

                is_valid, result = is_json_parseable(shadow_response.content)
                logger.info(f"Parseability status:\n{result}", )
                countd+=1
                logger.info(f"contd count is:\n{countd}",)

            logger.info("Success shadow repair!:",shadow_response.content)
            shadow_response = shadow_response.content
            logger.info(f"shadow_response type before: {type(shadow_response)}")
            logger.info(f"output type before: {type(output)}")

            shadow_response = json.loads(shadow_response)  # Convert JSON string into dictionary
            output = json.loads(output)  # Convert JSON string into dictionary

            logger.info(f"shadow_response type after: {type(shadow_response)} and output type after {type(output)}")

            output['edges']  = shadow_response['edges']  
            logger.info(f"{output} of {type(output)}")
            output = json.dumps(output, indent=2) # converts to dict to str
            is_valid_output, result = is_json_parseable(output)

            if is_valid_output == False:
                output = original_txt
                logger.info(f"The output was not parseable hence reverting to original_txt to this response:\n{output}")


        elif scenario == "gamified":

            if model_type=="gemini":
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_gamify_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))    
                    chain = PROMPTS.prompt_gamify_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges, llm=llm.bind(generation_config={"response_mime_type": "application/json"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges | llm.bind(generation_config={"response_mime_type": "application/json"})
            else:
                if repair_shadows_without_file != "1":
                    # chain = LLMChain(prompt=PROMPTS.prompt_gamify_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS.prompt_gamify_shadow_edges | llm.bind(response_format={"type": "json_object"})
                else:
                    # chain = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges, llm=llm.bind(response_format={"type": "json_object"}))
                    chain = PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges | llm.bind(response_format={"type": "json_object"})

            shadow_response = chain.invoke({"output": output,"language":language, "mpv":mpv, "mpv_string":mpv_string})
            is_valid, result = is_json_parseable(shadow_response.content)
            countd=0
            while not is_valid and countd<=3:
                txt = shadow_response.content
                modified_txt = re.findall(r'.*},', txt, re.DOTALL) # Finds last },
                if modified_txt:
                    modified_txt = modified_txt[0]  # Get the matched string
                else:
                    modified_txt = txt  # No match found, return original
                logger.info(f"original:::\n{txt}")
                logger.info(f"changed:::\n{modified_txt}")

                responses = modified_txt + "\n[CONTINUE_EXACTLY_FROM_HERE]" #changed txt
                logger.info(f"\nThe responses_modification to LLM is:\n{responses}",)

                if repair_shadows_without_file != "1":
                    # chain_edges_retry = LLMChain(prompt=PROMPTS.prompt_gamify_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS.prompt_gamify_shadow_edges_retry | llm
                else:
                    # chain_edges_retry = LLMChain(prompt=PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges_retry, llm=llm)
                    chain_edges_retry = PROMPTS_WITHOUT_FILE.prompt_gamify_shadow_edges_retry | llm

                response_retry = chain_edges_retry.invoke({"incomplete_response": responses, "output":output, "language":language, "mpv":mpv, "mpv_string":mpv_string})
                logger.info(f"response contd... is:\n{response_retry.content}",)

                responses = modified_txt + response_retry.content #changed modified_text to responses
                logger.info(f"responses+continued Combined is:\n{responses}",)
                
                shadow_response.content = responses

                is_valid, result = is_json_parseable(shadow_response.content)
                logger.info(f"Parseability status:\n{result}", )
                countd+=1
                logger.info(f"contd count is:\n{countd}",)

            logger.info("Success shadow repair!:",shadow_response.content)
            shadow_response = shadow_response.content
            logger.info(f"shadow_response type before: {type(shadow_response)}")
            logger.info(f"output type before: {type(output)}")

            shadow_response = json.loads(shadow_response)  # Convert JSON string into dictionary
            output = json.loads(output)  # Convert JSON string into dictionary

            logger.info(f"shadow_response type after: {type(shadow_response)} and output type after {type(output)}")

            output['edges']  = shadow_response['edges']  
            logger.info(f"{output} of {type(output)}")
            output = json.dumps(output, indent=2) # converts to dict to str
            is_valid_output, result = is_json_parseable(output)

            if is_valid_output == False:
                output = original_txt
                logger.info(f"The output was not parseable hence reverting to original_txt to this response:\n{output}")
 
        
    else:
        logger.info(f"Since error_flag is {error_flag}, no shadow edges found!")

    
    return output



def ANSWER_IMG(response_text, llm,relevant_doc,language,model_type):
    # prompt_template_img =PromptTemplate( 
    # input_variables=["response_text","context"],
    # template="""
    # Provided the context, look at the Images that are mentioned in the 'response_text': {response_text}. Provide a brief summary of those 
    # images stored in the 'context': {context}.
    # Format of Reply (The number of Images and their description may vary, depends on what is instructed in the
    # 'response_text'. If only one image is mentioned in the 'response_text', then you should include Image1 only. If there are 2 or more images then your reply should
    # also have same images as mentioned in the 'response_text'!):
    # {{"Image1": "file_name_..._page_..._image_...",
    # "Description1": "...",
    # "Image2": "file_name_..._page_..._image_...",
    # "Description2": "..."
    # and so on
    # }}
    # Warning: Include the complete schema of name defined. The complete schema of name includes
    # "file_name_..._page_..._image_..."
    # Take great care for the underscores. They are to be used exactly as defined. Also take 
    # extreme caution at the file_name since the file might be having its own - and _ which is not to be
    # tampered with in any way and should remain exactly the same!
    # [WARNING: The ... presents page number as int and image number as int. But, for the file_name_ it represents
    # as the file name itself which may have its own dashes or underscores or brackets. Whatever the file name
    # you found in the 'context', make sure you use the same name. ]
    # Answer():
    # """
    # )



    # chain = LLMChain(prompt=prompt_template_img,llm=llm)
    # img_response = chain.run({"response_text": response_text, "context": relevant_doc})
    # print("img_response is::",img_response)
###

    class image_loc(BaseModel):
        FileName: str = Field(description="Exact, absolutely Unchanged FileName of the image as mentioned in the 'Context'. FileName may contain special characters such as hyphens (-), underscores (_), semicolons (;), spaces, and others.")
        # PageNumber: Optional[str] = Field(description="If available, write PageNumber of the image. 'Null' if not available. !!!DO NOT USE PageNumber if SlideNumber is available.!!!")
        # SlideNumber: Optional[str] = Field(description="If available, slide number of the image.")
        # ImageNumber: int = Field(description="ImageNumber of the image")
        Description: str = Field(description="Description detail of the image")
        Logic: str = Field(description="Recommend out of MediaBlocks only (identify by title of pertinent MediaBlock) the pertinent image that shall be attached with. If an Image is not relevant to a MediaBlock, output 'NOT RELEVANT' as your response. The images with 'PIXMAPPED IMAGE' label are also 'NOT RELEVANT'.")

    class image(BaseModel):
        Image: List[image_loc] = Field(description="image_loc")

    parser = JsonOutputParser(pydantic_object=image)

    prompt = PromptTemplate(
    template="""
    You respond in the language of "{language}", since your responses are given to {language} speakers and they can only understand the language of {language}.
    
    !!!
    NEGATIVE PROMPT: RESPONDING OUTSIDE THE JSON FORMAT.   
    DO NOT START YOUR RESPONSE THAT STARTS WITH ```json and ENDS WITH ```
    The ticks you enclose JSON in, deems the JSON parseability unsuccessful.
    Just start the JSON response directly.
    Remove Line Breaks: Combined multi-line Description fields into single lines.
    Escape Single Quotes: Removed unnecessary escaping of single quotes since they're not needed in JSON strings.
    !!!    
    
    Search for those image or images only, whose descriptions in a MediaBlock of the 'Response Text' matches
    with the descriptions in the 'Context' data. Output only those image's or images' description from the 
    'Context' data.

    Most Important point to Note and be advised on it is that you should only select those images from the 
    'Context' whose description is relative to the MediaBlock description of the 'Response Text'. Your work will only be
    succesfull if you do implement this most important point. If an image description does not relative with
    the MediaBlock in the 'Response Text', then you respond in the format: 'NOT RELEVANT'. 
    The images with 'PIXMAPPED IMAGE' label are also 'NOT RELEVANT'.

    Please give logic in the Description section for why are you selecting an image in terms of relevance to
    the description of image in 'Context' to the description in the 'Response Text'.

    \n{format_instructions}\n'Response Text': {response_text}\n'Context': {context}""",
    input_variables=["response_text","context","language"],
    partial_variables={"format_instructions": parser.get_format_instructions()},
    )

    if model_type == "gemini":
        chain = prompt | llm | parser
    else:
        chain = prompt | llm | parser

    # just to debug the format_instructions and response_text
    format_instructions = parser.get_format_instructions()
    logger.info(f"response_text:\n{response_text}",)
    logger.info(f"format_instructions:\n{format_instructions}",)

    # invoking or running the img_response
    img_response = chain.invoke({"response_text": response_text, "context": relevant_doc, "language": language})
    logger.info(f"img_response is::{img_response}",)

    

###
    def create_structured_json(img_response):
        result = {}
        for index, img in enumerate(img_response['Image'], start=1):
            logger.info(f"img: {img}",)
            # if img['FileName']=="URL":
            #     # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
            #     image_key = f"FileName URL ImageNumber {img['ImageNumber']}"
            #     # Add the image key and description to the result dictionary
            #     result[f"Image{index}"] = image_key
            #     result[f"Description{index}"] = img['Description']
            #     result[f"Logic{index}"] = img['Logic']
            # else:
            #     if img['PageNumber'] is not None:
            #         # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
            #         image_key = f"FileName {img['FileName']} PageNumber {img['PageNumber']} ImageNumber {img['ImageNumber']}"
            #         # Add the image key and description to the result dictionary
            #         result[f"Image{index}"] = image_key
            #         result[f"Description{index}"] = img['Description']
            #         result[f"Logic{index}"] = img['Logic']
            #     else:
            #         # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
            #         image_key = f"FileName {img['FileName']} SlideNumber {img['SlideNumber']} ImageNumber {img['ImageNumber']}"
            #         # Add the image key and description to the result dictionary
            #         result[f"Image{index}"] = image_key
            #         result[f"Description{index}"] = img['Description']
            #         result[f"Logic{index}"] = img['Logic']

            # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
            image_key = f"{img['FileName']}"
            # Add the image key and description to the result dictionary
            result[f"Image{index}"] = image_key
            result[f"Description{index}"] = img['Description']
            result[f"Logic{index}"] = img['Logic']
                
        
        return json.dumps(result, indent=4)

    # Using the function to transform the data
    structured_response = create_structured_json(img_response)
    logger.info(structured_response)

    return str(structured_response)

#CODE

# import re

# text = "FileName vector-search-diagram- PageNumber:Null ImageNumber 2.JPG"

# pattern_FileName = r"FileName(.*?)\sPageNumber"


# pattern_PageNumber = r"PageNumber(.*?)\sImageNumber"

# pattern_SlideNumber = r"SlideNumber(.*?)\sImageNumber"

# pattern_ImageNumber = r"ImageNumber(.*?)\."


# match_FileName = re.search(pattern_PageNumber, text)
# if match_FileName:
#     filename = match_FileName.group(1)
#     print(filename)
# else:
#     print("No match found")
